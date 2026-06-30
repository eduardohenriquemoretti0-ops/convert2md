#!/usr/bin/env python3
"""Convert PDF, DOCX, XLSX, PPTX, HTML, CSV → Markdown."""
import argparse
import csv
import re
import sys
import time
from pathlib import Path

SUPPORTED = {".pdf", ".docx", ".xlsx", ".pptx", ".html", ".htm", ".csv"}


def _cell(v) -> str:
    if v is None:
        return ""
    return str(v).replace("\n", " ").replace("|", "\\|")


def _clean(text: str) -> str:
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines, out, prev = text.splitlines(), [], None
    for line in lines:
        line = line.rstrip()
        if line == "" and prev == "":
            continue
        out.append(line)
        prev = line
    return "\n".join(out).strip() + "\n"


def convert_pdf(path: Path) -> str:
    import pymupdf4llm
    return pymupdf4llm.to_markdown(str(path))


def convert_docx(path: Path) -> str:
    import mammoth
    import html2text as h2t
    with open(path, "rb") as f:
        result = mammoth.convert_to_html(f)
    h = h2t.HTML2Text()
    h.ignore_images = True
    h.body_width = 0
    return h.handle(result.value)


def convert_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for name in wb.sheetnames:
        ws = wb[name]
        rows = [r for r in ws.iter_rows(values_only=True) if any(c is not None for c in r)]
        if not rows:
            continue
        parts.append(f"## {name}\n")
        header = [_cell(c) for c in rows[0]]
        parts.append("| " + " | ".join(header) + " |")
        parts.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            cols = [_cell(c) for c in row]
            while len(cols) < len(header):
                cols.append("")
            parts.append("| " + " | ".join(cols) + " |")
        parts.append("")
    return "\n".join(parts)


def convert_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        title, body = "", []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            ph = getattr(shape, "placeholder_format", None)
            if ph is not None and ph.idx == 0:
                title = text
            else:
                body.append(text)
        parts.append(f"## Slide {i}{': ' + title if title else ''}")
        parts.extend(body)
        parts.append("")
    return "\n".join(parts)


def convert_html(path: Path) -> str:
    import html2text as h2t
    h = h2t.HTML2Text()
    h.ignore_images = True
    h.body_width = 0
    return h.handle(path.read_text(errors="replace"))


def convert_csv(path: Path) -> str:
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        rows = list(csv.reader(f))
    if not rows:
        return ""
    w = len(rows[0])
    parts = ["| " + " | ".join(rows[0]) + " |",
             "| " + " | ".join(["---"] * w) + " |"]
    for row in rows[1:]:
        row = (row + [""] * w)[:w]
        parts.append("| " + " | ".join(row) + " |")
    return "\n".join(parts) + "\n"


CONVERTERS = {
    ".pdf": convert_pdf,
    ".docx": convert_docx,
    ".xlsx": convert_xlsx,
    ".pptx": convert_pptx,
    ".html": convert_html,
    ".htm": convert_html,
    ".csv": convert_csv,
}


def convert(path: Path, out_dir: Path | None = None, stdout: bool = False) -> Path | None:
    ext = path.suffix.lower()
    fn = CONVERTERS.get(ext)
    if fn is None:
        print(f"skip: {path} (unsupported)", file=sys.stderr)
        return None
    try:
        md = _clean(fn(path))
    except Exception as e:
        print(f"error: {path}: {e}", file=sys.stderr)
        return None
    if stdout:
        sys.stdout.write(md)
        return None
    dest = (out_dir or path.parent) / (path.stem + ".md")
    dest.write_text(md, encoding="utf-8")
    print(f"{path.name} → {dest}")
    return dest


def watch(folder: Path, out_dir: Path | None = None) -> None:
    from watchdog.observers import Observer
    from watchdog.events import FileSystemEventHandler

    class Handler(FileSystemEventHandler):
        def _try(self, src: str) -> None:
            p = Path(src)
            if p.suffix.lower() in SUPPORTED and p.exists():
                convert(p, out_dir)

        def on_created(self, event):
            if not event.is_directory:
                self._try(event.src_path)

        def on_modified(self, event):
            if not event.is_directory:
                self._try(event.src_path)

    obs = Observer()
    obs.schedule(Handler(), str(folder), recursive=False)
    obs.start()
    print(f"watching {folder}  Ctrl+C to stop")
    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        obs.stop()
    obs.join()


def main() -> None:
    ap = argparse.ArgumentParser(description="Convert documents to Markdown")
    ap.add_argument("files", nargs="*", type=Path, metavar="FILE")
    ap.add_argument("-o", "--output-dir", type=Path, metavar="DIR")
    ap.add_argument("--stdout", action="store_true", help="write MD to stdout")
    ap.add_argument("--watch", type=Path, metavar="DIR", help="watch folder, auto-convert")
    args = ap.parse_args()

    if args.output_dir:
        args.output_dir.mkdir(parents=True, exist_ok=True)

    if args.watch:
        watch(args.watch, args.output_dir)
        return

    if not args.files:
        ap.print_help()
        sys.exit(1)

    for f in args.files:
        convert(f, args.output_dir, args.stdout)


if __name__ == "__main__":
    main()
